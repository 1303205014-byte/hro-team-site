#!/usr/bin/env python3
"""Generate HRO Team PPT from website content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)      # Deep navy
ACCENT = RGBColor(0x2C, 0x5F, 0xFF)       # Blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
MUTED = RGBColor(0x88, 0x88, 0x88)

def add_bg_shape(slide, color=DARK_BG):
    """Add full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei', line_spacing=1.3):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.line_spacing = line_spacing
    return txBox

def add_bullet_text(slide, left, top, width, height, items,
                    font_size=14, color=WHITE, bullet_color=ACCENT):
    """Add bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.line_spacing = 1.5
        p.space_before = Pt(4)
    return txBox

# ===== Slide 1: Title / Banner =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg_shape(slide1, DARK_BG)

# Accent line at top
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                prs.slide_width, Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Main title
add_text_box(slide1, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
             "打破人力瓶颈，释放组织增长力",
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "您的业务增长，是否正被人效低下、激励失灵或团队协同内耗所拖累？\n"
             "我们提供不止于方案的人力资源实战解方。",
             font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Stats bar
stats_y = Inches(4.8)
stats_data = [
    ("10+", "人均深耕年限"),
    ("1000+", "人规模运营"),
    ("3", "位核心专家"),
]
stat_width = Inches(3.5)
stat_gap = Inches(0.5)
start_x = (prs.slide_width - stat_width * 3 - stat_gap * 2) / 2

for i, (num, label) in enumerate(stats_data):
    x = start_x + i * (stat_width + stat_gap)
    add_text_box(slide1, x, stats_y, stat_width, Inches(0.6),
                 num, font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide1, x, stats_y + Inches(0.55), stat_width, Inches(0.4),
                 label, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Trust statement
add_text_box(slide1, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6),
             "我们的专家团队人均拥有超过10年的人力资源深耕经验，成功为超过1000人规模的企业设计并落地管理体系，"
             "并由3位在不同领域拥有标杆成果的核心专家领衔。",
             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Contact CTA
add_text_box(slide1, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
             "咨询热线：15014837925  |  专属对接人：边锋（Grace）",
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ===== Slide 2: Core Services =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide2, DARK_BG)

add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
             "核心服务", font_size=12, color=ACCENT, align=PP_ALIGN.LEFT)
add_text_box(slide2, Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.6),
             "我们专注于解决组织在发展中遇到的核心人力难题",
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

services = [
    {
        "title": "组织诊断与变革",
        "value": '解决"团队看似忙碌，却总推不动事、出不了活"的根源问题。',
        "items": [
            "组织架构诊断与效能分析",
            "岗位职责梳理与说明书编写",
            "人才盘点（九宫格/胜任力模型应用）",
            "关键岗位培训体系搭建",
        ]
    },
    {
        "title": "职位职级与薪酬体系搭建",
        "value": '解决"员工抱怨不公、骨干流失，薪酬成了成本而非激励工具"的体系问题。',
        "items": [
            "职级体系与宽幅薪酬设计",
            "员工晋升通道与发展路径规划",
            "薪酬套改方案设计与平稳落地",
            "长期激励（如股权、期权）机制设计",
        ]
    },
    {
        "title": "薪酬绩效改革与成本管控",
        "value": '解决"奖金发了没效果，人力成本连年涨，却不知花在哪、是否值得"的管控难题。',
        "items": [
            "销售/研发/管理等不同序列的绩效奖金与提成方案设计",
            "全面薪酬体系优化",
            "年度人力成本预算编制与动态管控机制",
            "基于市场数据的薪酬调研与竞争力分析",
        ]
    },
]

card_width = Inches(4)
card_height = Inches(5.2)
card_y = Inches(1.7)
card_gap = Inches(0.4)
start_x = (prs.slide_width - card_width * 3 - card_gap * 2) / 2

for i, svc in enumerate(services):
    x = start_x + i * (card_width + card_gap)
    # Card bg
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, card_y, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x36)
    card.line.fill.background()

    # Accent top bar
    bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, card_y, card_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    add_text_box(slide2, x + Inches(0.25), card_y + Inches(0.3),
                 card_width - Inches(0.5), Inches(0.5),
                 svc["title"], font_size=18, bold=True, color=WHITE)

    # Value statement
    add_text_box(slide2, x + Inches(0.25), card_y + Inches(0.85),
                 card_width - Inches(0.5), Inches(1.0),
                 svc["value"], font_size=13, color=LIGHT_GRAY)

    # Bullet items
    add_bullet_text(slide2, x + Inches(0.25), card_y + Inches(1.9),
                    card_width - Inches(0.5), Inches(3.0),
                    svc["items"], font_size=13, color=WHITE)

# ===== Slide 3: Expert Team - All 3 in one page =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide3, DARK_BG)

add_text_box(slide3, Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.4),
             "专家团队", font_size=12, color=ACCENT, align=PP_ALIGN.LEFT)
add_text_box(slide3, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.45),
             "历经千人体量企业的复杂管理实战，交付的不是理论，是能落地的结果",
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

experts = [
    {
        "name": "边锋  Grace",
        "badge": "薪酬数据与成本管控专家",
        "summary": "擅长将复杂的薪酬数据转化为清晰的决策依据，实现激励与成本的最优平衡。",
        "skills": "薪酬体系全模块设计、销售激励模型重塑、人力成本精细化预算与管控。",
        "results": '曾为一家快速扩张的互联网企业重塑销售激励制度，建立与业绩强挂钩的年度分配模型，有效牵引重点业务线增长超30%；主导搭建人力成本动态监测体系，通过数据化决策帮助企业在两年内实现人力成本费率下降5个百分点。',
    },
    {
        "name": "李文珍  Selina",
        "badge": "组织变革与高管激励实战派",
        "summary": "擅长从业务增长视角出发，设计能留住核心人才、激发团队战斗力的组织与激励方案。",
        "skills": "以业务为导向的组织诊断与变革、薪酬绩效体系顶层设计、核心人才长期激励规划。",
        "results": "主导某科技独角兽公司薪酬套改项目，针对性提升高价值员工薪酬竞争力，项目完成后关键岗位年度主动离职率降低40%；设计并落地高管差异化长期激励方案，成功保留全部核心管理层，并助力吸引多位行业高端人才加盟。",
    },
    {
        "name": "贾梅宝  Cecilia",
        "badge": "项目制组织与跨职能协同专家",
        "summary": '擅长破解矩阵式管理困局，通过组织重构与机制设计，让跨部门团队从"扯皮"变为"合力"。',
        "skills": "复杂业务下的组织架构设计、跨职能团队协同流程再造、项目制激励与考核机制。",
        "results": '针对某企业定制化业务丢单率高的问题，通过组织诊断锁定协同流程漏洞，主导设立"铁三角"项目制团队并绑定共同激励。在标杆项目中，将项目交付周期缩短20%，内部协作满意度提升35%，成功固化了一套可复制的项目制运营规则。',
    },
]

col_width = Inches(4)
col_height = Inches(5.8)
col_y = Inches(1.25)
col_gap = Inches(0.3)
start_x = (prs.slide_width - col_width * 3 - col_gap * 2) / 2

for i, exp in enumerate(experts):
    x = start_x + i * (col_width + col_gap)
    # Card bg
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, col_y, col_width, col_height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x36)
    card.line.fill.background()

    # Accent top bar
    bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, col_y, col_width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Photo placeholder (small, top)
    photo = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + Inches(0.2), col_y + Inches(0.25),
                                     Inches(1.2), Inches(1.2))
    photo.fill.solid()
    photo.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
    photo.line.fill.background()

    # Name + badge (to the right of photo)
    add_text_box(slide3, x + Inches(1.55), col_y + Inches(0.3), col_width - Inches(1.75), Inches(0.4),
                 exp["name"], font_size=16, bold=True, color=WHITE)
    add_text_box(slide3, x + Inches(1.55), col_y + Inches(0.7), col_width - Inches(1.75), Inches(0.35),
                 exp["badge"], font_size=11, color=ACCENT)

    # Summary
    add_text_box(slide3, x + Inches(0.2), col_y + Inches(1.6), col_width - Inches(0.4), Inches(0.9),
                 exp["summary"], font_size=11, color=LIGHT_GRAY)

    # Skills
    add_text_box(slide3, x + Inches(0.2), col_y + Inches(2.5), col_width - Inches(0.4), Inches(0.3),
                 "💡 核心擅长", font_size=11, bold=True, color=WHITE)
    add_text_box(slide3, x + Inches(0.2), col_y + Inches(2.8), col_width - Inches(0.4), Inches(0.9),
                 exp["skills"], font_size=10, color=LIGHT_GRAY)

    # Results
    add_text_box(slide3, x + Inches(0.2), col_y + Inches(3.7), col_width - Inches(0.4), Inches(0.3),
                 "✅ 代表成果", font_size=11, bold=True, color=WHITE)
    add_text_box(slide3, x + Inches(0.2), col_y + Inches(4.0), col_width - Inches(0.4), Inches(1.6),
                 exp["results"], font_size=10, color=LIGHT_GRAY)

# ===== Slide 4: Process =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide4, DARK_BG)

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
             "合作流程", font_size=12, color=ACCENT, align=PP_ALIGN.LEFT)
add_text_box(slide4, Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.6),
             "四步法专业服务，确保不是交付一份报告，而是真正解决你的问题",
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

steps = [
    ('01', '精准扫描，厘清根因',
     '你将获得：一份清晰的企业人力现状诊断报告，明确核心痛点与改进优先级，不再"凭感觉"管理。'),
    ('02', '量身定制，输出方案',
     '你将获得：一套紧密结合你业务特点的、可执行落地的详细解决方案与实施路径图。'),
     ('03', '嵌入业务，协同落地',
     '我们将：与你的人力及业务团队并肩工作，培训赋能、解答疑问，确保方案嵌入日常管理，平稳过渡。'),
    ('04', '复盘迭代，固化成效',
     '你将获得：项目效果的数据化复盘，明确投入产出比，并建立持续优化机制，让管理改进形成闭环。'),
]

step_width = Inches(5.8)
step_height = Inches(2.4)
positions = [
    (Inches(0.8), Inches(1.8)),
    (Inches(6.7), Inches(1.8)),
    (Inches(0.8), Inches(4.4)),
    (Inches(6.7), Inches(4.4)),
]

for i, (num, title, desc) in enumerate(steps):
    x, y = positions[i]
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, step_width, step_height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x36)
    card.line.fill.background()

    add_text_box(slide4, x + Inches(0.3), y + Inches(0.25), Inches(1), Inches(0.5),
                 num, font_size=28, bold=True, color=ACCENT)
    add_text_box(slide4, x + Inches(0.3), y + Inches(0.8), step_width - Inches(0.6), Inches(0.5),
                 title, font_size=18, bold=True, color=WHITE)
    add_text_box(slide4, x + Inches(0.3), y + Inches(1.3), step_width - Inches(0.6), Inches(0.9),
                 desc, font_size=14, color=LIGHT_GRAY)

# ===== Slide 5: Contact =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide5, DARK_BG)

add_text_box(slide5, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8),
             "还在为团队效能、人才流失或薪酬成本问题而困扰吗？",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
             "拖延的每一天，都可能意味着核心人才的动摇、业务机会的错失或不必要的成本消耗。",
             font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# CTA box
cta_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.5), Inches(4.0),
                                   Inches(4.3), Inches(1.2))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = ACCENT
cta_box.line.fill.background()
add_text_box(slide5, Inches(4.5), Inches(4.25), Inches(4.3), Inches(0.5),
             "立即行动，开启改变", font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(1), Inches(5.6), Inches(11.3), Inches(0.4),
             "咨询热线：15014837925  |  专属对接人：边锋（Grace）",
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(1), Inches(6.1), Inches(11.3), Inches(0.4),
             "扫码添加微信，直接与专家对话，获取针对性建议",
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
output_path = "/root/.openclaw/workspace/hro-team-site/HRO团队介绍.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
